from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Constants for styling
BACKGROUND_COLOR = RGBColor(0, 100, 0)
HEADER_RCB_RED = RGBColor(252, 5, 5)
HEADER_YELLOW = RGBColor(255, 255, 0)
HEADER_TEXT_WHITE = RGBColor(255, 255, 255)
QUESTION_TEXT_COLOR = RGBColor(255, 255, 0)
OPTION_TEXT_COLOR = RGBColor(255, 255, 255)
HEADER_BLUE = RGBColor(78, 157, 130)
FONT_NAME = "Calibri"
HEADER_FONT_SIZE = Pt(16)
QUESTION_FONT_SIZE = Pt(22)
OPTION_FONT_SIZE = Pt(20)

IMAGE_1_PATH = "logo.jpeg"
IMAGE_2_PATH = "testcase1.jpeg"

def create_custom_presentation(data):
    prs = Presentation()
    prs.slide_width = Inches(14)
    prs.slide_height = Inches(7.5)

    for i, question in enumerate(data["Question"]):
        options = data["Options"][i]
        year = None
        if "Year" in data:
            year = data["Year"][i] if i < len(data["Year"]) and data["Year"][i] else None
        add_custom_slide(prs, i + 1, question, options, year)

    prs.save(f"ppt_{datetime.now()}.pptx")

def add_custom_slide(prs, question_number, question, options, year):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Set slide background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    # Navigation Bar Sections (Placed in the same positions as before)
    nav_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0), Inches(5), Inches(0.5))
    nav_left.fill.solid()
    nav_left.fill.fore_color.rgb = HEADER_RCB_RED
    nav_left.text_frame.text = "NEXT LEVEL ACADEMY"
    nav_left.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_left.text_frame.paragraphs[0].font.bold = True
    nav_left.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
    nav_left.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_left.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    slide.shapes.add_picture(IMAGE_1_PATH, Inches(0), Inches(0), Inches(1), Inches(0.5))

    nav_middle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(0), Inches(5), Inches(0.5))
    nav_middle.fill.solid()
    nav_middle.fill.fore_color.rgb = HEADER_YELLOW
    nav_middle.text_frame.text = "Indian Geography Introduction of India"
    nav_middle.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_middle.text_frame.paragraphs[0].font.bold = True
    nav_middle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    nav_middle.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    slide.shapes.add_picture(IMAGE_2_PATH, Inches(5.5), Inches(0), Inches(1), Inches(0.5))

    nav_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), Inches(0), Inches(3), Inches(0.5))
    nav_right.fill.solid()
    nav_right.fill.fore_color.rgb = HEADER_RCB_RED
    nav_right.text_frame.text = "BY: Chaman"
    nav_right.text_frame.paragraphs[0].font.size = HEADER_FONT_SIZE
    nav_right.text_frame.paragraphs[0].font.bold = True
    nav_right.text_frame.paragraphs[0].font.color.rgb = HEADER_TEXT_WHITE
    nav_right.text_frame.paragraphs[0].font.name = FONT_NAME
    nav_right.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Adjust content to the upper right side
    horizontal_offset = Inches(8)  # Move to the right side (adjust as needed)
    vertical_offset_question_number = Inches(1)  # Move the question number lower toward the bottom (adjusted)

    # Add a small red background behind the question number (move to right upper side)
    question_number_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, horizontal_offset, vertical_offset_question_number, Inches(0.4), Inches(0.7))
    question_number_bg.fill.solid()
    question_number_bg.fill.fore_color.rgb = HEADER_RCB_RED

    # Add Question Number (Centered in the box, White Text)
    question_number_box = slide.shapes.add_textbox(horizontal_offset, vertical_offset_question_number, Inches(0.4), Inches(0.7))
    question_number_frame = question_number_box.text_frame
    question_number_frame.clear()
    question_number_frame.paragraphs[0].text = f"{question_number}."  # Display question number
    question_number_frame.paragraphs[0].font.size = QUESTION_FONT_SIZE
    question_number_frame.paragraphs[0].font.bold = True
    question_number_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    question_number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centered horizontally
    question_number_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centered vertically

    # Add Question Text (Now horizontally aligned with question number)
    question_box = slide.shapes.add_textbox(horizontal_offset + Inches(0.5), Inches(0.8), Inches(4), Inches(1.5))
    question_frame = question_box.text_frame
    question_frame.clear()
    question_frame.word_wrap = True
    question_paragraph = question_frame.add_paragraph()
    question_paragraph.text = question  # Question text without number
    question_paragraph.font.size = QUESTION_FONT_SIZE
    question_paragraph.font.bold = True
    question_paragraph.font.color.rgb = QUESTION_TEXT_COLOR
    question_paragraph.font.name = FONT_NAME
    question_paragraph.alignment = PP_ALIGN.LEFT
    question_paragraph.line_spacing = 1

    # Dynamically adjust the height of the question box
    question_lines = len(question) // 40 + 1  # Estimate lines based on text length
    question_box.height = Inches(0.5 + question_lines * 0.5)

    # Check if question_box height exceeds 50% of slide height (7.5 inches)
    if question_box.height > Inches(3.75):  # 50% of 7.5 inches slide height
        # Reduce font size and adjust the height of the question box
        question_paragraph.font.size = Pt(16)  # Reduce font size
        question_box.height = Inches(3)  # Set a smaller height to avoid overlap

    # Add "Year Mentioned" if year is provided
    if year:
        year_top =question_box.height + 2 + Inches(0.4)  # Add some spacing
        year_text = f"RRB NTPC {year}."
        year_box = slide.shapes.add_textbox(Inches(10), year_top, Inches(12), Inches(0.5))
        year_frame = year_box.text_frame
        year_frame.clear()
        year_paragraph = year_frame.add_paragraph()
        year_paragraph.text = year_text
        year_paragraph.font.size = QUESTION_FONT_SIZE
        year_paragraph.font.bold = False
        year_paragraph.font.color.rgb = RGBColor(0, 255, 255)
        year_paragraph.font.name = FONT_NAME
        year_paragraph.alignment = PP_ALIGN.LEFT
        options_top=year_box.top + year_box.height + Inches(0.1)
    else:
        options_top=question_box.top + question_box.height + Inches(0.2)

    # Adjust vertical alignment of the options with added space between question and options
    options_box = slide.shapes.add_textbox(Inches(8.5), options_top, Inches(12), Inches(3))
    options_frame = options_box.text_frame
    options_frame.clear()
    for option in options:
        option_paragraph = options_frame.add_paragraph()
        option_paragraph.text = option
        option_paragraph.font.size = OPTION_FONT_SIZE
        option_paragraph.font.color.rgb = OPTION_TEXT_COLOR
        option_paragraph.font.name = FONT_NAME
        option_paragraph.alignment = PP_ALIGN.LEFT
        option_paragraph.space_after = Pt(8)

new_data = {
    "Question": [
        "जंतु विज्ञान (Zoology) अध्ययन करता है:",
        "सूची-I का सूची-II से सुमेल कीजिए और नीचे दिए कूट का प्रयोग करते हुए सही उत्तर का चयन कीजिए:",
        "फूलों के अध्ययन को क्या कहते हैं?",
        "कांटों के वैज्ञानिक अध्ययन को क्या कहते हैं?",
        "हमारे शरीर में आनुवंशिकता की इकाई को कहते हैं:",
        "निम्न में से किस कोशिकांग (Organelle) को कोशिका का 'पावर हाउस' कहते हैं?",
        "रेशम कीट पालन को क्या कहते हैं?",
        "प्लाज़्मा झिल्ली बनी होती है:",
        "'विटीकल्चर' के द्वारा निम्नलिखित में से कौन एक उत्पादित होता है?",
        "आनुवंशिकी निम्न में से किसका अध्ययन है?"
    ],
    "Options": [
        ["(a) केवल जीवित जानवरों का", "(b) केवल जीवित वनस्पति का", "(c) जीवित व मृत जानवरों दोनों का", "(d) जीवित व मृत वनस्पति दोनों का"],
        [
            "सूची-I: \n(A) पक्षी \n(B) वंशागति \n(C) पर्यावरण \n(D) जीवाश्म",
            "सूची-II: \n1. पेलियोजूलॉजी \n2. इकोलॉजी \n3. ऑर्निथोलॉजी \n4. जेनेटिक्स",
            "कूट: \n(a) A-1, B-3, C-4, D-2",
            "(b) A-3, B-4, C-2, D-1",
            "(c) A-4, B-2, C-1, D-3",
            "(d) A-2, B-4, C-1, D-3"
        ],
        ["(a) फ्लोरीकलॉजी", "(b) एंथोलॉजी", "(c) एग्रोस्टोलॉजी", "(d) पेलिनोलॉजी"],
        ["(a) इक्थियोलॉजी", "(b) एंटोमोलॉजी", "(c) पैरासिटोलॉजी", "(d) मेलाकोलॉजी"],
        ["(a) गुणसूत्र", "(b) डीएनए", "(c) जीन", "(d) केन्द्रक"],
        ["(a) राइबोसोम", "(b) माइटोकॉन्ड्रिया", "(c) कोशिका झिल्ली", "(d) सेंट्रोसोम"],
        ["(a) एपिकल्चर", "(b) हॉर्टीकल्चर", "(c) सेरीकल्चर", "(d) फ्लोरीकल्चर"],
        ["(a) प्रोटीन से", "(b) लिपिड से", "(c) कार्बोहाइड्रेट से", "(d) दोनों (a) तथा (b)"],
        ["(a) सिल्क", "(b) केंचुआ", "(c) शहर", "(d) अंगूर"],
        ["(a) मेंडल का नियम", "(b) जीव विकास", "(c) डी.एन.ए. संरचना", "(d) आनुवांशिकता और वितरण"]
    ],
    "Answers": [
        "(c)",
        "(b)",
        "(b)",
        "(b)",
        "(c)",
        "(b)",
        "(c)",
        "(d)",
        "(d)",
        "(d)"
    ],
    "Year": [
        "U.P.P.C.S. (Pre)",
        "U.P. Lower Sub. (Pre) 1998",
        "Jharkhand P.C.S. (Pre) 2003",
        "Uttarakhand U.D.A./L.D.A. (Pre) 2003",
        "U.P.P.S.C. (GIC) 2010",
        "U.P.R.O./A.R.O. (Mains) 2013",
        "Uttarakhand U.D.A./L.D.A. (Pre) 2003",
        "U.P.P.C.S. (Mains) 2008",
        "U.P.P.C.S. (Mains) 2003",
        "53rd to 55th B.P.S.C. (Pre) 2011"
    ]
}


# Create the presentation
create_custom_presentation(new_data)
